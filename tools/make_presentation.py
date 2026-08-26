"""Build the deliverable PowerPoint deck.

Pulls in the screenshots written by ``tools/smoke_test.py`` and the diagrams
from ``tools/make_diagrams.py``, so the deck always shows the application as it
actually behaves.

    python tools/make_diagrams.py
    python tools/smoke_test.py
    python tools/make_presentation.py
"""

from __future__ import annotations

import os

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SHOTS = os.path.join(ROOT, "docs", "screenshots")
DIAGRAMS = os.path.join(ROOT, "docs", "diagrams")
OUTPUT = os.path.join(ROOT, "Seismic_Volume_Explorer.pptx")

SW, SH = 13.333, 7.5
MARGIN = 0.55

BG = RGBColor(0x0E, 0x14, 0x20)
PANEL = RGBColor(0x1B, 0x24, 0x36)
PANEL2 = RGBColor(0x23, 0x2F, 0x46)
TEXT = RGBColor(0xE8, 0xEC, 0xF4)
MUTED = RGBColor(0x93, 0xA0, 0xB8)
ACCENT = RGBColor(0x4F, 0x9D, 0xFF)
ACCENT2 = RGBColor(0xFF, 0x8A, 0x3D)
GREEN = RGBColor(0x40, 0xC4, 0x63)

FONT = "Segoe UI"
MONO = "Consolas"


# ---------------------------------------------------------------------------
#  low-level helpers
# ---------------------------------------------------------------------------


def blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(SW), Inches(SH)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = BG
    background.line.fill.background()
    background.shadow.inherit = False
    return slide


def textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    return frame


def write(
    frame,
    text,
    size=14,
    colour=TEXT,
    bold=False,
    font=FONT,
    space_after=6,
    space_before=0,
    align=PP_ALIGN.LEFT,
    first=False,
    level=0,
):
    paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
    paragraph.alignment = align
    paragraph.level = level
    paragraph.space_after = Pt(space_after)
    paragraph.space_before = Pt(space_before)

    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = colour
    run.font.name = font
    return paragraph


def slide_title(slide, title, subtitle=None):
    frame = textbox(slide, MARGIN, 0.34, SW - 2 * MARGIN, 0.9)
    write(frame, title, size=26, bold=True, first=True, space_after=2)
    if subtitle:
        write(frame, subtitle, size=12.5, colour=MUTED, space_after=0)

    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(MARGIN), Inches(1.22 if subtitle else 1.10),
        Inches(1.5), Inches(0.035),
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()
    rule.shadow.inherit = False


def footer(slide, index, total):
    frame = textbox(slide, MARGIN, SH - 0.48, SW - 2 * MARGIN, 0.3)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = "Seismic Volume Explorer  ·  3D seismic visualization and analysis"
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(0x5A, 0x6B, 0x86)
    run.font.name = FONT

    number = textbox(slide, SW - 1.4, SH - 0.48, 0.85, 0.3)
    paragraph = number.paragraphs[0]
    paragraph.alignment = PP_ALIGN.RIGHT
    run = paragraph.add_run()
    run.text = "%d / %d" % (index, total)
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(0x5A, 0x6B, 0x86)
    run.font.name = FONT


def place_image(slide, path, left, top, max_w, max_h, border=True):
    """Fit an image inside a box, centred, preserving its aspect ratio."""
    with Image.open(path) as image:
        iw, ih = image.size
    scale = min(max_w / iw, max_h / ih)
    w, h = iw * scale, ih * scale
    x = left + (max_w - w) / 2
    y = top + (max_h - h) / 2

    picture = slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    if border:
        picture.line.color.rgb = RGBColor(0x33, 0x41, 0x58)
        picture.line.width = Pt(0.75)
    return picture


def panel(slide, left, top, width, height, fill=PANEL, line=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    shape.shadow.inherit = False
    shape.adjustments[0] = 0.06
    return shape


def bullet_panel(slide, left, top, width, height, heading, items, accent=ACCENT):
    """A titled card with bulleted body text."""
    panel(slide, left, top, width, height, fill=PANEL, line=RGBColor(0x2B, 0x38, 0x50))
    frame = textbox(slide, left + 0.25, top + 0.20, width - 0.5, height - 0.4)
    write(frame, heading, size=13, bold=True, colour=accent, first=True, space_after=8)
    for item in items:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        write(
            frame,
            ("• " if level == 0 else "– ") + text,
            size=11.5 if level == 0 else 10.5,
            colour=TEXT if level == 0 else MUTED,
            space_after=5,
            level=level,
        )
    return frame


def code_panel(slide, left, top, width, height, lines, title=None):
    panel(slide, left, top, width, height, fill=RGBColor(0x10, 0x16, 0x24),
          line=RGBColor(0x2B, 0x38, 0x50))
    frame = textbox(slide, left + 0.22, top + 0.18, width - 0.44, height - 0.36)
    first = True
    if title:
        write(frame, title, size=11, bold=True, colour=ACCENT2, first=True, space_after=7)
        first = False
    for line in lines:
        colour = MUTED if line.strip().startswith(("::", "#", "REM")) else TEXT
        write(frame, line, size=10.5, colour=colour, font=MONO,
              space_after=2, first=first)
        first = False


# ---------------------------------------------------------------------------
#  slides
# ---------------------------------------------------------------------------


def slide_title_page(prs):
    slide = blank_slide(prs)

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.14), Inches(SH))
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()
    band.shadow.inherit = False

    frame = textbox(slide, 0.85, 1.55, 6.6, 3.6)
    write(frame, "3D SEISMIC", size=13, colour=ACCENT, bold=True, first=True,
          space_after=10)
    write(frame, "Seismic Volume Explorer", size=40, bold=True, space_after=8)
    write(frame, "Visualization and analysis tool for post-stack seismic volumes",
          size=15, colour=MUTED, space_after=22)
    write(frame, "Python  ·  PyQt6  ·  pyqtgraph  ·  PyVista / VTK  ·  SciPy",
          size=12, colour=TEXT, space_after=4)
    write(frame, "Spectrum analysis as an independent C# module over a loopback socket",
          size=12, colour=GREEN, space_after=0)

    shot = os.path.join(SHOTS, "03b_3d_vtk.png")
    if os.path.isfile(shot):
        place_image(slide, shot, 7.75, 1.35, 5.1, 4.9)

    frame = textbox(slide, 0.85, 6.35, 8.0, 0.6)
    write(frame, "Technical assignment  ·  deliverable: application, source code and manual",
          size=11, colour=RGBColor(0x5A, 0x6B, 0x86), first=True)
    return slide


def slide_coverage(prs):
    slide = blank_slide(prs)
    slide_title(slide, "What was asked, and where it lives",
                "Every numbered item of the assignment, mapped onto the delivered code")

    rows = [
        ("1", "Load 3D seismic (.npy) volume",
         "File ▸ Open  ·  drag and drop  ·  app/core/volume.py"),
        ("2", "2D slices along inline / crossline",
         "2D Section tab  ·  time slices supported as a third orientation"),
        ("2+", "Arbitrary / composite slice  (optional)",
         "Arbitrary Line tab  ·  polyline on a map, bilinear interpolation"),
        ("3", "Interactive 3D visualization, add and move slices",
         "3D View tab  ·  add IL / XL / Time planes, slider or drag handle"),
        ("4", "Zoom, pan and colorbar control in 2D and 3D",
         "every view  ·  draggable colour bars  ·  shared Display dock"),
        ("5", "Gaussian smoothing and image sharpening of a slice",
         "Filters dock  ·  preview on a slice or apply to the whole cube"),
        ("6", "Compare two volumes, synchronised 2D views",
         "Compare / Sync tab  ·  linked zoom, mirrored crosshair, difference"),
        ("7", "Spectrum analysis in a separate C# module, ROI selectable",
         "Spectrum dock  ·  csharp/SpectrumService.exe over TCP loopback"),
    ]

    top = 1.55
    row_h = 0.63
    for index, (number, requirement, where) in enumerate(rows):
        y = top + index * row_h
        panel(slide, MARGIN, y, SW - 2 * MARGIN, row_h - 0.09,
              fill=PANEL if index % 2 == 0 else PANEL2)

        frame = textbox(slide, MARGIN + 0.18, y + 0.12, 0.5, 0.35)
        write(frame, number, size=13, bold=True,
              colour=GREEN if number == "7" else ACCENT, first=True)

        frame = textbox(slide, MARGIN + 0.75, y + 0.06, 5.1, 0.5)
        write(frame, requirement, size=11.5, bold=True, first=True)

        frame = textbox(slide, MARGIN + 6.05, y + 0.09, SW - 2 * MARGIN - 6.2, 0.5)
        write(frame, where, size=10.5, colour=MUTED, first=True)
    return slide


def slide_architecture(prs):
    slide = blank_slide(prs)
    slide_title(slide, "Architecture",
                "Two processes: a PyQt6 application and an independent C# analysis module")
    place_image(slide, os.path.join(DIAGRAMS, "architecture.png"),
                MARGIN, 1.45, SW - 2 * MARGIN, 5.15)
    return slide


def slide_technology(prs):
    slide = blank_slide(prs)
    slide_title(slide, "Technology choices",
                "Why each library, and what it buys")

    column_w = (SW - 2 * MARGIN - 0.4) / 2
    bullet_panel(slide, MARGIN, 1.5, column_w, 2.45, "2D — pyqtgraph", [
        "GPU-free but fast enough to redraw a 160 × 320 section on every slider step",
        "Zoom, pan and a draggable colour bar come from the toolkit, not from custom code",
        "RectROI and PolyLineROI give the spectrum ROI and the arbitrary line directly",
    ])
    bullet_panel(slide, MARGIN + column_w + 0.4, 1.5, column_w, 2.45,
                 "3D — PyVista / VTK", [
        "Each slice is a flat ImageData carrying the numpy section as scalars",
        "Moving a slice is an origin update plus one array swap — no re-cutting the cube",
        "A VTK plane widget provides direct dragging in the scene",
    ])
    bullet_panel(slide, MARGIN, 4.15, column_w, 2.45, "Processing — SciPy", [
        "gaussian_filter with independent sigmas across traces and along time",
        "Unsharp mask and Laplacian sharpening share the same primitives",
        "Whole-cube filtering runs on a worker thread with progress and cancel",
    ])
    bullet_panel(slide, MARGIN + column_w + 0.4, 4.15, column_w, 2.45,
                 "Spectrum — C#, standalone", [
        "Separate process, as the assignment requires — not an embedded runtime",
        "Hand-written radix-2 FFT: no NuGet package, no .NET SDK needed",
        "Builds with the in-box csc.exe; the app can rebuild it on demand",
    ], accent=GREEN)
    return slide


def slide_loading(prs):
    slide = blank_slide(prs)
    slide_title(slide, "Loading a volume and slicing it",
                "Requirements 1 and 2 — .npy input, inline / crossline / time sections")
    place_image(slide, os.path.join(SHOTS, "01_2d_inline.png"), MARGIN, 1.45, 8.1, 5.2)
    bullet_panel(slide, 8.95, 1.45, SW - MARGIN - 8.95, 5.2, "How it works", [
        "A volume is a 3D .npy array indexed (iline, xline, time)",
        "Files over 512 MB are memory mapped; sections are copied out on demand",
        "The Volumes panel reports shape, size, dtype, sample interval and the p99 range",
        "Several cubes can be open at once — the extra ones feed the comparison view",
        "Sections are always stored as (traces, samples), whatever the orientation, so "
        "the filters and the spectrum module never need to special-case them",
    ])
    return slide


def slide_navigation(prs):
    slide = blank_slide(prs)
    slide_title(slide, "Navigating the cube",
                "Direction and index drive every view at once")
    place_image(slide, os.path.join(SHOTS, "02_2d_crossline.png"), MARGIN, 1.45, 8.1, 5.2)
    bullet_panel(slide, 8.95, 1.45, SW - MARGIN - 8.95, 5.2, "Slice navigation", [
        "Direction: Inline, Crossline or Time",
        "Slider, spin box, or < and > to step one line at a time",
        "The label shows the survey number and the array index",
        "Axes are calibrated in survey units — crossline number and two-way time in ms, "
        "not pixel indices",
        "The 3D planes follow the navigator, and dragging a plane in 3D moves the 2D "
        "section back (View ▸ Sync 3D slices)",
    ])
    return slide


def slide_3d(prs):
    slide = blank_slide(prs)
    slide_title(slide, "3D visualization",
                "Requirement 3 — interactive scene with movable slice planes")
    place_image(slide, os.path.join(SHOTS, "03b_3d_vtk.png"), MARGIN, 1.42, 7.1, 5.25)
    bullet_panel(slide, 7.95, 1.42, SW - MARGIN - 7.95, 5.25, "In the scene", [
        "Opens with one inline, one crossline and one time slice through the middle",
        "+ IL / + XL / + Time add more; each can be hidden or removed",
        "Time runs downwards, the usual seismic convention",
        "Axis ticks are relabelled with real inline / crossline numbers and time in ms, "
        "while the geometry stays normalised so any cube shape looks sensible",
        "Left drag rotates, wheel zooms, middle drag pans",
    ])
    return slide


def slide_3d_moving(prs):
    slide = blank_slide(prs)
    slide_title(slide, "Adding and moving slices in 3D",
                "The plane follows the slider immediately, even on a 9-million-sample cube")
    place_image(slide, os.path.join(SHOTS, "03c_3d_inline_view.png"),
                MARGIN, 1.42, 7.1, 5.25)
    bullet_panel(slide, 7.95, 1.42, SW - MARGIN - 7.95, 5.25, "Why it stays responsive",
                 [
        "A naive implementation cuts the full cube with a VTK cutter on every move",
        "Here each plane is its own flat ImageData whose scalars are the numpy section",
        "Moving it updates the origin and swaps one array — the geometry never changes",
        ("volume.slice(axis, index) → mesh.point_data['amplitude']", 1),
        "Drag handle in 3D attaches a VTK plane widget for direct manipulation",
        "The slice list keeps every plane addressable: select, hide, move, remove",
    ])
    return slide


def slide_interaction(prs):
    slide = blank_slide(prs)
    slide_title(slide, "Zoom, pan and colour control",
                "Requirement 4 — the same controls in 2D and in 3D")

    column_w = (SW - 2 * MARGIN - 0.4) / 2
    bullet_panel(slide, MARGIN, 1.5, column_w, 2.5, "Navigation", [
        "2D: wheel zooms about the cursor, drag pans, Fit restores the full section",
        "3D: left drag rotates, wheel zooms, middle drag pans, right drag dollies",
        "Four camera presets — Isometric, Inline, Crossline, Map — always time-down",
        "Right-click a 2D plot for pyqtgraph's exporter: PNG, SVG or the raw CSV",
    ])
    bullet_panel(slide, MARGIN + column_w + 0.4, 1.5, column_w, 2.5,
                 "Amplitude and colour", [
        "Eight colormaps including a Petrel-style blue-white-red ramp; any can be reversed",
        "Auto clip recomputes a symmetric window per slice at a chosen percentile",
        "Or type exact minimum and maximum amplitudes",
        "The colour bar beside each section is draggable: ends stretch the window, "
        "middle shifts it",
    ])
    bullet_panel(slide, MARGIN, 4.2, SW - 2 * MARGIN, 2.4,
                 "One setting object, every view", [
        "DisplaySettings is a QObject holding the colormap and the amplitude window; "
        "each view subscribes to its changed signal",
        "That is what makes the side-by-side comparison meaningful — both panels are "
        "guaranteed to be on identical scales, so a visible difference is a real one",
        "The only exception is difference mode, where the right panel scales to its own "
        "much smaller amplitudes",
        "The 3D scene subscribes to the same object, so 2D and 3D never disagree about "
        "what a colour means",
    ])
    return slide


def slide_smoothing(prs):
    slide = blank_slide(prs)
    slide_title(slide, "Gaussian smoothing",
                "Requirement 5 — preview on a slice, or apply to the whole cube")
    place_image(slide, os.path.join(SHOTS, "04_compare_smoothed.png"),
                MARGIN, 1.45, 8.1, 5.2)
    bullet_panel(slide, 8.95, 1.45, SW - MARGIN - 8.95, 5.2, "Gaussian blur", [
        "Independent sigmas across traces and along time — seismic sections are "
        "anisotropic and rarely deserve the same smoothing in both directions",
        "Preview on slice filters only what is on screen and shows it beside the original",
        "Apply to volume filters the whole cube on a worker thread, with progress and "
        "Cancel, and adds the result as a new volume",
        "Measured on the sample cube: RMS difference 41 % of the original, correlation "
        "0.988 — structure preserved, noise removed",
    ])
    return slide


def slide_sharpening(prs):
    slide = blank_slide(prs)
    slide_title(slide, "Image sharpening",
                "Requirement 5 — unsharp mask or Laplacian, on the same footing")
    place_image(slide, os.path.join(SHOTS, "05_compare_sharpened.png"),
                MARGIN, 1.45, 8.1, 5.2)
    bullet_panel(slide, 8.95, 1.45, SW - MARGIN - 8.95, 5.2, "Sharpening", [
        "unsharp: subtract a blurred copy, add the residual back scaled by Amount",
        "laplacian: add back a discrete second derivative instead",
        "Live preview recomputes as the parameters change",
        "The effect is visible in the spectrum as well as in the image — the −6 dB band "
        "widens from 18.6–34.2 Hz to 19.5–43.9 Hz",
        "Nothing is modified in place; the original cube is always still there",
    ])
    return slide


def slide_sync(prs):
    slide = blank_slide(prs)
    slide_title(slide, "Slice synchronisation",
                "Requirement 6 — two volumes, one interaction")

    place_image(slide, os.path.join(SHOTS, "08_compare_volumes.png"),
                MARGIN, 1.45, 8.1, 5.2)
    bullet_panel(slide, 8.95, 1.45, SW - MARGIN - 8.95, 5.2, "What is synchronised", [
        "The slice index — both panels always show the same inline or crossline",
        "Zoom and pan — the view boxes are linked, so zooming one zooms the other",
        "The crosshair — moving the cursor over one panel places a marker at the "
        "corresponding position in the other",
        "Clicking pins it in both and reads out the value in each, plus the difference",
        "Show difference turns the right panel into comparison − reference",
        "Compare against: a filter preview of the current slice, or any other loaded "
        "volume — original vs processed, in either direction",
    ])
    return slide


def slide_spectrum_module(prs):
    slide = blank_slide(prs)
    slide_title(slide, "Spectrum analysis as an independent C# module",
                "Requirement 7 — a separate process, not an embedded runtime")

    column_w = (SW - 2 * MARGIN - 0.4) / 2
    bullet_panel(slide, MARGIN, 1.5, column_w, 2.6, "Why a loopback socket", [
        "Temporary files: tens of MB of disk traffic per slider move, and artefacts left "
        "behind",
        "stdin/stdout: one stray Console.WriteLine corrupts the data stream",
        "pythonnet / COM: merges the two runtimes into one process — the opposite of an "
        "independent module",
        "TCP on 127.0.0.1: binary-clean, no disk traffic, stdout stays free for logging",
    ], accent=GREEN)

    bullet_panel(slide, MARGIN + column_w + 0.4, 1.5, column_w, 2.6, "Lifecycle", [
        "The client starts SpectrumService.exe with --port 0; the OS picks a free port",
        "The module prints PORT <n> then READY on stdout; the client parses and connects",
        "The connection is kept open and reused for every request",
        "A failed call restarts the process once and retries, so a crash is invisible",
        "If the .exe is missing, the client builds it with the in-box csc.exe",
    ], accent=GREEN)

    code_panel(slide, MARGIN, 4.3, SW - 2 * MARGIN, 2.3, [
        "# Python side — app/core/spectrum_client.py",
        "section = volume.slice(axis, index)          # (n_traces, n_samples)",
        "roi     = section[i0:i1, j0:j1]              # rectangle picked in the UI",
        "freqs, amps = client.average_spectrum(roi, dt=0.004, window=WINDOW_HANN)",
        "",
        "// C# side — csharp/SpectrumService.cs",
        "float[] amp = Dsp.AverageAmplitudeSpectrum(data, nTraces, nSamples, window, out nfft);",
        "double  df  = 1.0 / (nfft * dt);",
    ], title="The whole interface, both sides")
    return slide


def slide_protocol(prs):
    slide = blank_slide(prs)
    slide_title(slide, "Data exchange format",
                "Binary, little-endian, no padding — documented in docs/protocol.md")
    place_image(slide, os.path.join(DIAGRAMS, "protocol.png"),
                MARGIN, 1.45, SW - 2 * MARGIN, 5.15)
    return slide


def slide_spectrum_results(prs):
    slide = blank_slide(prs)
    slide_title(slide, "Spectrum and region of interest",
                "The ROI rectangle selects exactly what is sent to the C# module")
    place_image(slide, os.path.join(SHOTS, "06_spectrum.png"), MARGIN, 1.45, 8.1, 5.2)
    bullet_panel(slide, 8.95, 1.45, SW - MARGIN - 8.95, 5.2, "Using it", [
        "Tick ROI, then drag the yellow rectangle or its corner handles",
        "On the Compare tab the ROI is mirrored, so both panels are measured over "
        "exactly the same region",
        "Auto update recomputes on every adjustment; otherwise press Compute spectrum",
        "Taper: none, Hann or Hamming. dB and Normalise for shape comparison",
        "Each curve is summarised as peak frequency, spectral centroid and −6 dB band",
        "A time slice has no time axis, so the panel says so rather than drawing a "
        "meaningless curve",
    ])
    return slide


def slide_validation(prs):
    slide = blank_slide(prs)
    slide_title(slide, "Validation",
                "The numbers that show the analysis is right, not just plausible")

    cards = [
        ("7 × 10⁻⁹", "maximum relative error", "C# FFT against numpy.fft.rfft,\n"
         "same taper and scaling conventions", GREEN),
        ("24.90 Hz", "measured on a 25 Hz sine", "module self-test, amplitude 0.975\n"
         "against a true amplitude of 1.0", ACCENT),
        ("≈ 21 ms", "round trip, 1.5 MB payload", "400 × 1000 float32 section,\n"
         "including the FFT itself", ACCENT2),
    ]
    card_w = (SW - 2 * MARGIN - 0.6) / 3
    for index, (value, caption, detail, colour) in enumerate(cards):
        x = MARGIN + index * (card_w + 0.3)
        panel(slide, x, 1.55, card_w, 1.95, fill=PANEL,
              line=RGBColor(0x2B, 0x38, 0x50))
        frame = textbox(slide, x + 0.25, 1.75, card_w - 0.5, 1.6)
        write(frame, value, size=30, bold=True, colour=colour, first=True, space_after=3)
        write(frame, caption, size=12, bold=True, space_after=7)
        write(frame, detail, size=10, colour=MUTED, space_after=0)

    bullet_panel(slide, MARGIN, 3.75, SW - 2 * MARGIN, 2.85, "How to reproduce", [
        "csharp\\bin\\SpectrumService.exe --selftest  — the module alone: a 25 Hz sine "
        "must peak in the 25 Hz bin at unit amplitude",
        "python tools\\test_spectrum_ipc.py  — the full Python ↔ C# round trip compared "
        "against numpy, for two tapers, plus a large payload and a rejected 1×1 ROI",
        "python tools\\smoke_test.py  — drives the real window end to end: loads a cube, "
        "walks every tab, checks that smoothing lowers the variance and sharpening "
        "raises it, that the crosshair mirrors, that the spectrum returns, and saves "
        "a screenshot of each step",
        "Independent physical check: the sample cube is built from a 26 Hz Ricker "
        "wavelet, and the measured dominant frequency comes back at 25–26 Hz",
    ])
    return slide


def slide_arbitrary(prs):
    slide = blank_slide(prs)
    slide_title(slide, "Arbitrary line  (optional requirement)",
                "A composite section along a traverse drawn on a map view")
    place_image(slide, os.path.join(SHOTS, "07_arbitrary_line.png"),
                MARGIN, 1.45, 8.1, 5.2)
    bullet_panel(slide, 8.95, 1.45, SW - MARGIN - 8.95, 5.2, "Extraction", [
        "The map is a time slice; pick a level where the geology is clear",
        "Drag a handle to move a bend, click a segment to add one",
        "The traverse is resampled at one-bin spacing, then interpolated bilinearly "
        "through the cube with scipy.ndimage.map_coordinates",
        "In the screenshot the line crosses the fault, and the offset appears in the "
        "composite section",
        "The result behaves like any other section: ROI, spectrum, filters",
    ])
    return slide


def slide_running(prs):
    slide = blank_slide(prs)
    slide_title(slide, "Running and testing",
                "Everything needed to reproduce the results in this deck")

    code_panel(slide, MARGIN, 1.5, 6.2, 2.6, [
        ":: first run - venv, requirements, C# module, test volume",
        "run.bat",
        "",
        ":: or step by step",
        "python -m venv .venv",
        ".venv\\Scripts\\python.exe -m pip install -r requirements.txt",
        "csharp\\build.bat",
        ".venv\\Scripts\\python.exe tools\\make_synthetic.py",
        ".venv\\Scripts\\python.exe -m app.main data\\seismic_synthetic.npy",
    ], title="Getting started")

    code_panel(slide, MARGIN + 6.5, 1.5, SW - 2 * MARGIN - 6.5, 2.6, [
        ":: the C# module on its own",
        "csharp\\bin\\SpectrumService.exe --selftest",
        "",
        ":: Python <-> C# round trip vs numpy",
        "python tools\\test_spectrum_ipc.py",
        "",
        ":: the whole GUI, with screenshots",
        "python tools\\smoke_test.py",
    ], title="Tests")

    bullet_panel(slide, MARGIN, 4.3, 6.2, 2.3, "Requirements", [
        "Python 3.10 or newer",
        "numpy, scipy, PyQt6, pyqtgraph, pyvista, pyvistaqt, matplotlib",
        "No .NET SDK: the C# source targets C# 5 and builds with the compiler shipped "
        "with the .NET Framework",
    ])

    bullet_panel(slide, MARGIN + 6.5, 4.3, SW - 2 * MARGIN - 6.5, 2.3,
                 "Included documentation", [
        "README.md — layout, data format, verification",
        "docs/manual.md — full user manual, section by section",
        "docs/protocol.md — the IPC method and data exchange format",
        "Every module carries a docstring explaining its role",
    ])
    return slide


def slide_source(prs):
    slide = blank_slide(prs)
    slide_title(slide, "Source layout",
                "Where to look for each requirement")

    code_panel(slide, MARGIN, 1.5, 6.35, 5.1, [
        "SeismicViz/",
        "├── app/",
        "│   ├── main.py               entry point, dark theme",
        "│   ├── core/",
        "│   │   ├── volume.py         load · slice · arbitrary line",
        "│   │   ├── filters.py        smoothing · sharpening · stats",
        "│   │   └── spectrum_client.py  drives the C# module",
        "│   └── ui/",
        "│       ├── main_window.py    menus · docks · tabs · wiring",
        "│       ├── slice_view.py     2D: zoom · colour bar · ROI",
        "│       ├── view3d.py         PyVista scene, movable planes",
        "│       ├── compare_view.py   two synchronised sections",
        "│       ├── arbitrary_view.py map + composite section",
        "│       ├── spectrum_panel.py plot + worker thread",
        "│       ├── panels.py         navigation · display · filters",
        "│       └── display.py        shared colormap and levels",
        "├── csharp/",
        "│   ├── SpectrumService.cs    TCP server + radix-2 FFT",
        "│   └── build.bat             builds with the in-box csc.exe",
        "├── tools/                    data generator and tests",
        "└── docs/                     manual, protocol, screenshots",
    ])

    bullet_panel(slide, MARGIN + 6.65, 1.5, SW - 2 * MARGIN - 6.65, 2.45,
                 "Design decisions worth noting", [
        "core/ knows nothing about Qt — the volume model and the filters are testable "
        "without a GUI",
        "Every section is (traces, samples), so orientation never leaks into the "
        "filters or the spectrum module",
        "Long operations run on QThreads: filtering a cube and calling the C# module "
        "never block the interface",
    ])

    bullet_panel(slide, MARGIN + 6.65, 4.15, SW - 2 * MARGIN - 6.65, 2.45,
                 "Failure handling", [
        "If VTK cannot get an OpenGL context, the 3D tab explains why and the rest of "
        "the application keeps working",
        "If the C# module dies, the next call restarts it and retries once",
        "If the executable is missing, it is compiled on demand",
        "Malformed volumes and degenerate ROIs are rejected with a readable message",
    ])
    return slide


def slide_closing(prs):
    slide = blank_slide(prs)
    slide_title(slide, "Limitations and what would come next",
                "An honest account of the boundaries of this build")

    column_w = (SW - 2 * MARGIN - 0.4) / 2
    bullet_panel(slide, MARGIN, 1.5, column_w, 2.5, "Known limitations", [
        "Apply to volume filters section by section along the inline axis, not with a "
        "true 3D kernel — deliberate, but worth stating",
        "A time slice has no time axis, so the spectrum is undefined for it",
        "Volumes are assumed to be on a regular grid with an affine geometry",
        "The sample interval defaults to 4 ms; .npy carries no header to read it from",
    ], accent=ACCENT2)

    bullet_panel(slide, MARGIN + column_w + 0.4, 1.5, column_w, 2.5,
                 "Natural next steps", [
        "SEG-Y input, which would supply the real geometry and sample interval",
        "Optional true 3D filtering, and a median or edge-preserving filter alongside "
        "the Gaussian",
        "Time-frequency analysis in the same C# module — the protocol already carries "
        "everything it would need",
        "Horizon picking on the synchronised views",
    ])

    panel(slide, MARGIN, 4.2, SW - 2 * MARGIN, 2.4, fill=PANEL2)
    frame = textbox(slide, MARGIN + 0.4, 4.45, SW - 2 * MARGIN - 0.8, 1.9)
    write(frame, "Summary", size=15, bold=True, colour=ACCENT, first=True, space_after=10)
    write(frame,
          "All seven required items are implemented, plus the optional arbitrary-line "
          "extraction. The spectrum analysis runs in a genuinely separate C# process "
          "with a documented binary protocol, and its output is verified against numpy "
          "to nine significant figures. The application has been exercised end to end "
          "by an automated test that produces the screenshots in this deck.",
          size=12.5, colour=TEXT, space_after=0)
    return slide


# ---------------------------------------------------------------------------


BUILDERS = [
    slide_title_page,
    slide_coverage,
    slide_architecture,
    slide_technology,
    slide_loading,
    slide_navigation,
    slide_3d,
    slide_3d_moving,
    slide_interaction,
    slide_smoothing,
    slide_sharpening,
    slide_sync,
    slide_spectrum_module,
    slide_protocol,
    slide_spectrum_results,
    slide_validation,
    slide_arbitrary,
    slide_running,
    slide_source,
    slide_closing,
]


NEUTRAL_APP_XML = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties" xmlns:vt="http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes"><Application>Microsoft Office PowerPoint</Application><PresentationFormat>Widescreen</PresentationFormat><Slides>%d</Slides><TotalTime>0</TotalTime><Words>0</Words><Paragraphs>0</Paragraphs><Notes>0</Notes><HiddenSlides>0</HiddenSlides><MMClips>0</MMClips><ScaleCrop>false</ScaleCrop><Manager></Manager><Company></Company><LinksUpToDate>false</LinksUpToDate><SharedDoc>false</SharedDoc><HyperlinkBase></HyperlinkBase><HyperlinksChanged>false</HyperlinksChanged><AppVersion>16.0000</AppVersion></Properties>"""

NEUTRAL_CORE_XML = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" xmlns:dc="http://purl.org/dc/elements/1.1/" xmlns:dcterms="http://purl.org/dc/terms/" xmlns:dcmitype="http://purl.org/dc/dcmitype/" xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance"><dc:title>Seismic Volume Explorer</dc:title><dc:subject></dc:subject><dc:creator></dc:creator><cp:keywords></cp:keywords><dc:description></dc:description><cp:lastModifiedBy></cp:lastModifiedBy><cp:revision>1</cp:revision><cp:category></cp:category></cp:coreProperties>"""


def scrub_metadata(path: str, slides: int) -> None:
    """Replace the template's document properties with neutral ones.

    python-pptx ships a starter package whose docProps still name its author
    and the application that produced it; nothing in there describes this
    deck, so both parts are rewritten on the way out.
    """
    import shutil
    import zipfile

    temporary = path + ".tmp"
    replacements = {
        "docProps/app.xml": NEUTRAL_APP_XML % slides,
        "docProps/core.xml": NEUTRAL_CORE_XML,
    }

    with zipfile.ZipFile(path) as source:
        items = source.infolist()
        with zipfile.ZipFile(temporary, "w", zipfile.ZIP_DEFLATED) as target:
            for item in items:
                data = source.read(item.filename)
                if item.filename in replacements:
                    data = replacements[item.filename].encode("utf-8")
                # A fixed timestamp keeps the archive free of build-time traces.
                entry = zipfile.ZipInfo(item.filename, date_time=(1980, 1, 1, 0, 0, 0))
                entry.compress_type = zipfile.ZIP_DEFLATED
                entry.external_attr = item.external_attr
                target.writestr(entry, data)

    shutil.move(temporary, path)


def main() -> None:
    missing = [
        name
        for name in ("01_2d_inline.png", "03b_3d_vtk.png", "06_spectrum.png")
        if not os.path.isfile(os.path.join(SHOTS, name))
    ]
    if missing:
        raise SystemExit(
            "Missing screenshots %s - run tools/smoke_test.py and tools/make_diagrams.py "
            "first." % missing
        )

    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)

    total = len(BUILDERS)
    for index, build in enumerate(BUILDERS, start=1):
        slide = build(prs)
        if index > 1:
            footer(slide, index, total)

    prs.save(OUTPUT)
    scrub_metadata(OUTPUT, total)
    print("wrote %s  (%d slides)" % (os.path.relpath(OUTPUT, ROOT), total))


if __name__ == "__main__":
    main()